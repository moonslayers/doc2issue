#!/usr/bin/env python3
"""Extrae texto y renderiza slides como imágenes PNG.

Usa LibreOffice si está disponible para convertir PPTX → PDF → imágenes.
Si no, usa fallback: extrae imágenes incrustadas con contexto de slide.

Uso:
    uv run python3 scripts/extract_pptx.py docs/archivo.pptx [output/]
"""
import sys, json, shutil, subprocess, zipfile
from pathlib import Path
from pptx import Presentation


def extract_via_libreoffice(pptx: Path, out: Path, pptx_stem: str) -> tuple[list[str], list[str]]:
    """Renderiza slides como imágenes vía LibreOffice + PyMuPDF."""
    import fitz  # pymupdf (import aquí porque es opcional para el fallback)

    pdf_path = out / f"{pptx_stem}.pdf"
    subprocess.run([
        "libreoffice", "--headless", "--convert-to", "pdf",
        "--outdir", str(out), str(pptx),
    ], check=True, capture_output=True)

    images = []
    doc = fitz.open(pdf_path)
    for i, page in enumerate(doc, 1):
        pix = page.get_pixmap(dpi=200)
        img_name = f"{pptx_stem}_slide_{i:03d}.png"
        img_path = out / "images" / img_name
        pix.save(str(img_path))
        pix = None
        images.append(str(img_path))
    doc.close()
    pdf_path.unlink()  # limpiar PDF temporal
    return images


def extract_fallback(pptx: Path, out: Path, pptx_stem: str, prs) -> tuple[list[str], list[str]]:
    """Fallback: imágenes incrustadas agrupadas por slide (sin renderizado)."""
    import fitz  # solo para crear una imagen básica si no hay imágenes incrustadas
    del fitz  # no lo usamos realmente aquí

    images = []
    # Extraer imágenes incrustadas del ZIP del PPTX
    with zipfile.ZipFile(pptx) as z:
        for name in z.namelist():
            if name.startswith("ppt/media/"):
                img_name = f"{pptx_stem}_{Path(name).name}"
                img_path = out / "images" / img_name
                if not img_path.exists():
                    img_path.write_bytes(z.read(name))
                images.append(str(img_path))

    print("⚠️  LibreOffice no disponible. Usando extracción básica de imágenes.",
          file=sys.stderr)
    print("💡 Para renderizar slides completos: sudo pacman -S libreoffice-fresh",
          file=sys.stderr)
    return images


def extract(pptx_path: str, output_dir: str = "output"):
    pptx = Path(pptx_path)
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    (out / "images").mkdir(exist_ok=True)

    prs = Presentation(pptx)
    slides_text = []

    # Extraer texto slide por slide
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        slides_text.append(f"## Slide {i}\n" + "\n".join(texts) +
                          (f"\n\n**Notas:** {notes}" if notes else ""))

    # Renderizar o fallback
    has_libreoffice = shutil.which("libreoffice") is not None
    if has_libreoffice:
        images = extract_via_libreoffice(pptx, out, pptx.stem)
    else:
        images = extract_fallback(pptx, out, pptx.stem, prs)

    md_path = out / f"{pptx.stem}.md"
    md_path.write_text("\n\n---\n\n".join(slides_text), encoding="utf-8")

    manifest = {
        "source": str(pptx),
        "type": "pptx",
        "markdown_file": str(md_path),
        "slides": len(prs.slides),
        "images": images,
        "image_type": "full_slide" if has_libreoffice else "embedded",
    }
    (out / f"{pptx.stem}.manifest.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    print(json.dumps(manifest, indent=2))


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: uv run python3 scripts/extract_pptx.py docs/archivo.pptx [output/]",
              file=sys.stderr)
        sys.exit(1)
    extract(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else "output")
