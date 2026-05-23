"""Fixtures compartidos para tests de scripts de extracción.

Genera documentos mínimos de prueba para cada formato.
Cada fixture crea un archivo temporal y lo limpia al finalizar.
"""
import pytest
from pathlib import Path


@pytest.fixture
def test_pdf(tmp_path):
    """Crea un PDF mínimo con texto e imagen incrustada."""
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Requerimiento: Login con Google OAuth", fontsize=12)
    page.insert_text((50, 80), "Prioridad: Alta", fontsize=10)
    page.insert_text((50, 110), "Criterios de aceptación:", fontsize=10)

    # Crear una imagen simple (rectángulo rojo) para incrustar
    pix = fitz.Pixmap(fitz.csRGB, 100, 100, True)
    pix.set_rect(pix.irect, (255, 0, 0))
    img_bytes = pix.tobytes("png")

    page.insert_image(fitz.Rect(50, 140, 150, 240), stream=img_bytes)

    pdf_path = tmp_path / "test.pdf"
    doc.save(str(pdf_path))
    doc.close()
    return pdf_path


@pytest.fixture
def test_docx(tmp_path):
    """Crea un DOCX mínimo con texto e imágenes."""
    from docx import Document
    from docx.shared import Inches
    doc = Document()
    doc.add_heading("Especificación Módulo de Pagos", 1)
    doc.add_paragraph(
        "El módulo debe permitir pagos con tarjeta de crédito y débito."
    )
    doc.add_paragraph("Criterios de aceptación:", style="List Bullet")
    doc.add_paragraph("Validar número de tarjeta", style="List Bullet")
    doc.add_paragraph("Soportar VISA y Mastercard", style="List Bullet")

    docx_path = tmp_path / "test.docx"
    doc.save(str(docx_path))
    return docx_path


@pytest.fixture
def test_pptx(tmp_path):
    """Crea un PPTX mínimo con texto y notas del presentador."""
    from pptx import Presentation
    prs = Presentation()

    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide1.shapes.title.text = "Nuevo Dashboard"
    txBox = slide1.shapes.add_textbox(left=100, top=100, width=200, height=50)
    txBox.text_frame.text = "Gráficos en tiempo real"

    # Slide 2 con notas
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Arquitectura Propuesta"
    notes_slide = slide2.notes_slide
    notes_slide.notes_text_frame.text = "Usar microservicios con cola de eventos"

    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def test_xlsx(tmp_path):
    """Crea un XLSX mínimo con columnas tipo backlog."""
    import pandas as pd
    df = pd.DataFrame({
        "Title": ["Login OAuth", "Dashboard", "Reportes PDF"],
        "Description": [
            "Implementar login con Google",
            "Dashboard con métricas en tiempo real",
            "Exportar reportes a PDF",
        ],
        "Priority": ["High", "Medium", "Low"],
        "Size": ["M", "L", "S"],
        "Estimate Hours": [8, 16, 4],
    })
    xlsx_path = tmp_path / "test.xlsx"
    df.to_excel(str(xlsx_path), index=False)
    return xlsx_path
